import streamlit as st
import io, re, zipfile, traceback
import fitz  # PyMuPDF
import docx  # python-docx
import pptx  # python-pptx
import pandas as pd
import pytesseract
from PIL import Image
import numpy as np
import faiss
from sentence_transformers import SentenceTransformer
from unidecode import unidecode
from openai import OpenAI
from fpdf import FPDF  # Pour générer le PDF

# --- CLIENT API NVIDIA ---
def get_client():
    return OpenAI(
        base_url="https://integrate.api.nvidia.com/v1",
        api_key="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"  # votre clé ici
    )

# --- FONCTIONS D'EXTRACTION (inchangées) ---
def extract_text_from_pdf(file_obj):
    file_obj.seek(0)
    text = ""
    try:
        doc = fitz.open(stream=file_obj.read(), filetype="pdf")
        for page in doc:
            text += page.get_text() + "\n"
    except Exception as e:
        st.error(f"Erreur extraction PDF: {e}")
    return text

def extract_text_from_docx(file_obj):
    file_obj.seek(0)
    text = ""
    try:
        doc = docx.Document(file_obj)
        for para in doc.paragraphs:
            text += para.text + "\n"
    except Exception as e:
        st.error(f"Erreur extraction DOCX: {e}")
    return text

def extract_text_from_pptx(file_obj):
    file_obj.seek(0)
    text = ""
    try:
        prs = pptx.Presentation(file_obj)
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    text += shape.text + "\n"
    except Exception as e:
        st.error(f"Erreur extraction PPTX: {e}")
    return text

def extract_text_from_xlsx(file_obj):
    file_obj.seek(0)
    text = ""
    try:
        df = pd.read_excel(file_obj)
        text = df.to_string()
    except Exception as e:
        st.error(f"Erreur extraction XLSX: {e}")
    return text

def extract_text_from_txt(file_obj):
    file_obj.seek(0)
    try:
        return file_obj.read().decode('utf-8')
    except Exception as e:
        st.error(f"Erreur extraction TXT: {e}")
        return ''

def extract_text_from_image(file_obj):
    file_obj.seek(0)
    text = ''
    try:
        img = Image.open(file_obj)
        text = pytesseract.image_to_string(img)
    except Exception as e:
        st.error(f"Erreur extraction image: {e}")
    return text

def extract_text_from_zip(file_obj):
    file_obj.seek(0)
    full = ''
    try:
        with zipfile.ZipFile(file_obj) as z:
            for name in z.namelist():
                ext = name.split('.')[-1].lower()
                with z.open(name) as f:
                    data = io.BytesIO(f.read())
                    if ext == 'pdf':
                        full += extract_text_from_pdf(data)
                    elif ext == 'docx':
                        full += extract_text_from_docx(data)
                    elif ext == 'pptx':
                        full += extract_text_from_pptx(data)
                    elif ext == 'xlsx':
                        full += extract_text_from_xlsx(data)
                    elif ext == 'txt':
                        full += extract_text_from_txt(data)
                    elif ext in ['jpg','jpeg','png']:
                        full += extract_text_from_image(data)
    except Exception as e:
        st.error(f"Erreur extraction ZIP: {e}")
    return full

def process_file(file_obj, filename):
    ext = filename.split('.')[-1].lower()
    if ext == 'zip':       return extract_text_from_zip(file_obj)
    if ext == 'pdf':       return extract_text_from_pdf(file_obj)
    if ext == 'docx':      return extract_text_from_docx(file_obj)
    if ext == 'pptx':      return extract_text_from_pptx(file_obj)
    if ext == 'xlsx':      return extract_text_from_xlsx(file_obj)
    if ext == 'txt':       return extract_text_from_txt(file_obj)
    if ext in ['jpg','jpeg','png']:
        return extract_text_from_image(file_obj)
    return ''

# --- NETTOYAGE & CHUNKING ---
def clean_text(text):
    t = text.lower()
    t = unidecode(t)
    t = re.sub(r"[^a-z0-9\s\.,;:?!'-]", ' ', t)
    return re.sub(r'\s+', ' ', t)

def split_text_into_chunks(text, chunk_size=500, overlap=50):
    words = text.split()
    chunks = []
    i = 0
    while i < len(words):
        chunks.append(' '.join(words[i:i+chunk_size]))
        i += chunk_size - overlap
    return chunks

# --- EMBEDDING & INDEX ---
def embed_chunks(chunks, model):
    return model.encode(chunks, convert_to_numpy=True)

def build_faiss_index(emb):
    idx = faiss.IndexFlatL2(emb.shape[1])
    idx.add(emb)
    return idx

def search_best_chunks(index, qemb, k=8):
    _, ids = index.search(np.array([qemb]), k)
    return ids[0]

# --- PROMPT DYNAMIQUE ---
def generate_dynamic_prompt(context, sections, language='francais'):
    if language == 'francais':
        header = (
            'Tu es un assistant spécialisé en recherche scientifique. '
            'Génère un résumé structuré et sous forme de paragraphes avec :\n'
            '- Introduction\n'
        )
    else:
        header = (
            'You are a scientific assistant. '
            'Generate a structured summary in paragraph form with:\n'
            '- Introduction\n'
        )
    for s in sections:
        header += f'- {s}\n'
    header += (
        '- Conclusion\n\n'
        'Contenu extrait :\n' + context +
        '\n\nRédige le résumé suivant cette structure.'
    )
    return header

# --- OUTPUT PDF/DOCX ---
def generate_pdf(text):
    pdf = FPDF()
    pdf.add_page()
    pdf.set_font('Arial', size=12)
    for line in text.split('\n'):
        pdf.multi_cell(0, 8, unidecode(line))
    return pdf.output(dest='S').encode('latin1')

def generate_docx(text):
    doc = docx.Document()
    for line in text.split('\n'):
        doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue()

# --- INTERFACE STREAMLIT ---
st.title('Résumé automatique de documents scientifiques')

# Initialise le résumé en session
if 'resume_text' not in st.session_state:
    st.session_state['resume_text'] = None

# Options
ing = st.selectbox('Langue du résumé', ['Français','Anglais'])
lang = 'francais' if ing=='Français' else 'english'
preset = ['Contexte et objectifs','Verrous','Approche globale','Principaux résultats','Discussion','Perspectives']
chosen = st.multiselect('Sections à inclure', preset, default=['Contexte et objectifs','Principaux résultats','Perspectives'])
txt = st.text_input('Autres sections (séparées par virgules)','')
custom = [s.strip() for s in txt.split(',') if s.strip()]
sections = chosen + [s for s in custom if s not in chosen]

mode = st.radio('Mode', ['Fichier unique','Dossier (.zip)'])
files = st.file_uploader(
    'Upload',
    type=['pdf','docx','pptx','xlsx','txt','jpg','jpeg','png','zip'],
    accept_multiple_files=(mode.startswith('Dossier'))
)

# Génération du résumé
if st.button('Générer le résumé'):
    if not files:
        st.warning('Veuillez uploader au moins un fichier.')
        st.stop()

    uploaded = files if isinstance(files, list) else [files]
    txts = []

    for f in uploaded:
        if f is None:
            st.error('Fichier non valide détecté.')
            st.stop()
        # On lit tout en mémoire
        data = f.read()
        buffer = io.BytesIO(data)
        try:
            txt = process_file(buffer, f.name)
        except Exception as e:
            st.error(f"Erreur lors du traitement de {f.name}: {e}")
            st.error(traceback.format_exc())
            st.stop()
        if not txt:
            st.warning(f"Aucun texte extrait de {f.name}.")
        txts.append(txt)

    full = "\n".join(txts)
    if not full.strip():
        st.error('Aucun texte extrait.')
        st.stop()

    clean = clean_text(full)
    chks = split_text_into_chunks(clean)

    with st.spinner('Embeddings et recherche...'):
        mdl = SentenceTransformer('all-MiniLM-L6-v2')
        emb = embed_chunks(chks, mdl)
        idx = build_faiss_index(np.array(emb))
        qemb = mdl.encode(
            'objectif verrous méthodologie résultats perspectives',
            convert_to_numpy=True
        )
        ids = search_best_chunks(idx, qemb)
        ctx = "\n\n".join(chks[i] for i in ids if i < len(chks))
        prm = generate_dynamic_prompt(ctx, sections, lang)
        res = get_client().chat.completions.create(
            model='mistralai/mistral-small-24b-instruct',
            messages=[{'role':'user','content':prm}],
            temperature=0.2, top_p=0.7, max_tokens=2048, stream=True
        )
        out = ''
        for chunk in res:
            d = chunk.choices[0].delta.content
            if d:
                out += d

    st.success('Résumé terminé.')
    st.session_state['resume_text'] = out

# Affichage + téléchargement si résumé dispo
if st.session_state['resume_text']:
    out = st.session_state['resume_text']
    st.text_area('Résumé généré', out, height=300)

    download_format = st.radio(
        'Format de téléchargement', ['PDF','DOCX'], key='fmt_choice'
    )
    if download_format == 'PDF':
        st.download_button(
            'Télécharger en PDF',
            data=generate_pdf(out),
            file_name='resume.pdf',
            mime='application/pdf'
        )
    else:
        st.download_button(
            'Télécharger en DOCX',
            data=generate_docx(out),
            file_name='resume.docx',
            mime='application/vnd.openxmlformats-officedocument.wordprocessingml.document'
        )
